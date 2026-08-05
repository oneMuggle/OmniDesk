"""smart_assistant/extractors/office_extractor.py — 统一 Office 文件抽取器

按扩展名路由到各格式处理器（python-docx / mammoth / pdfplumber /
openpyxl / python-pptx / pandas），输出结构化文本 + 表格 + sheet 信息，
供聊天附件上下文注入与工具读取。所有抽取失败统一抛 OfficeExtractError。
"""

from __future__ import annotations

import io
from dataclasses import dataclass, field
from pathlib import Path

import pandas as pd


class OfficeExtractError(Exception):
    """Office 文件抽取失败（格式不支持 / 损坏 / 超限）。"""


@dataclass
class ExtractedDocument:
    """一次抽取的产物。text 为段落合并文本；markdown 优先 docx/pdf；sheets 供 xlsx。"""

    text: str = ""
    markdown: str = ""
    tables: list = field(default_factory=list)
    sheets: list = field(default_factory=list)
    metadata: dict = field(default_factory=dict)
    format: str = ""


class OfficeExtractor:
    """按扩展名路由抽取，并提供切片与 prompt 注入工具。"""

    MAX_UPLOAD_SIZE = 10 * 1024 * 1024
    CHUNK_SIZE = 8000
    INLINE_LIMIT = 50_000
    SUPPORTED_EXTENSIONS = {".docx", ".pdf", ".xlsx", ".pptx", ".txt", ".md", ".csv"}

    @staticmethod
    def extract(file) -> ExtractedDocument:
        name = getattr(file, "name", "") or ""
        ext = Path(name).suffix.lower()
        if ext not in OfficeExtractor.SUPPORTED_EXTENSIONS:
            raise OfficeExtractError(
                f"暂不支持 {ext or '该'} 格式，支持 .docx/.pdf/.xlsx/.pptx/.txt/.md/.csv"
            )
        # 优先用 file.size 预检大小，避免超大文件先整读入内存
        size_hint = getattr(file, "size", None)
        if size_hint is not None and size_hint > OfficeExtractor.MAX_UPLOAD_SIZE:
            raise OfficeExtractError("文件超过 10MB 上限")
        try:
            data = file.read()
        except Exception as exc:  # pragma: no cover — 依赖具体文件对象
            raise OfficeExtractError(f"读取文件失败: {exc}") from exc
        finally:
            if hasattr(file, "seek"):
                file.seek(0)

        if len(data) > OfficeExtractor.MAX_UPLOAD_SIZE:
            raise OfficeExtractError("文件超过 10MB 上限")

        handler = {
            ".docx": OfficeExtractor._extract_docx,
            ".pdf": OfficeExtractor._extract_pdf,
            ".xlsx": OfficeExtractor._extract_xlsx,
            ".pptx": OfficeExtractor._extract_pptx,
            ".txt": OfficeExtractor._extract_text,
            ".md": OfficeExtractor._extract_text,
            ".csv": OfficeExtractor._extract_csv,
        }[ext]

        try:
            doc = handler(io.BytesIO(data), name)
        except OfficeExtractError:
            raise
        except Exception as exc:
            raise OfficeExtractError("文件无法解析，请确认文件未损坏") from exc

        doc.metadata.update({"filename": name, "size": len(data), "format": ext.lstrip(".")})
        return doc

    @staticmethod
    def _extract_docx(file, name) -> ExtractedDocument:
        from docx import Document
        import mammoth

        document = Document(file)
        paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
        tables = []
        for table in document.tables:
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            if rows:
                tables.append({"headers": rows[0], "rows": rows[1:]})
        file.seek(0)
        md_result = mammoth.convert_to_markdown(file)
        markdown = md_result.value
        return ExtractedDocument(
            text="\n".join(paragraphs),
            markdown=markdown,
            tables=tables,
            format="docx",
        )

    @staticmethod
    def _extract_pdf(file, name) -> ExtractedDocument:
        import pdfplumber

        tables = []
        pages_text = []
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                text = page.extract_text()  # 昂贵操作，单次调用并缓存
                if text:
                    pages_text.append(text)
                for t in page.extract_tables() or []:
                    rows = [[cell if cell is not None else "" for cell in row] for row in t]
                    if rows:
                        tables.append({"headers": rows[0], "rows": rows[1:]})
        return ExtractedDocument(text="\n".join(pages_text), tables=tables, format="pdf")

    @staticmethod
    def _extract_xlsx(file, name) -> ExtractedDocument:
        from openpyxl import load_workbook

        wb = load_workbook(file, read_only=True, data_only=True)
        sheets, parts = [], []
        for ws in wb.worksheets:
            rows = [list(r) for r in ws.iter_rows(values_only=True) if any(r)]
            if not rows:
                continue
            headers = [str(c) if c is not None else "" for c in rows[0]]
            data = [[str(c) if c is not None else "" for c in r] for r in rows[1:]]
            sheets.append({"name": ws.title, "headers": headers, "data": data})
            parts.append(
                f"### Sheet: {ws.title}\n"
                + pd.DataFrame(data, columns=headers).to_markdown(index=False)
            )
        wb.close()
        return ExtractedDocument(
            text="\n\n".join(parts),
            markdown="\n\n".join(parts),
            sheets=sheets,
            format="xlsx",
        )

    @staticmethod
    def _extract_pptx(file, name) -> ExtractedDocument:
        from pptx import Presentation

        prs = Presentation(file)
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text)
                if getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        texts.append(" | ".join(cell.text for cell in row.cells))
            parts.append(f"--- 第 {i} 页 ---\n" + "\n".join(t for t in texts if t.strip()))
        return ExtractedDocument(text="\n".join(parts), format="pptx")

    @staticmethod
    def _extract_text(file, name) -> ExtractedDocument:
        raw = file.read().decode("utf-8", errors="replace")
        return ExtractedDocument(text=raw, format=Path(name).suffix.lstrip("."))

    @staticmethod
    def _extract_csv(file, name) -> ExtractedDocument:
        # 编码回退：优先 utf-8，中文 GBK/GB18030 CSV 在 utf-8 下会 UnicodeDecodeError
        raw = file.read()
        df = None
        for encoding in ("utf-8", "gbk", "gb18030"):
            try:
                df = pd.read_csv(io.BytesIO(raw), encoding=encoding)
                break
            except UnicodeDecodeError:
                continue
        if df is None:
            raise OfficeExtractError("CSV 编码无法识别，仅支持 utf-8/gbk/gb18030")
        return ExtractedDocument(
            text=df.to_markdown(index=False),
            markdown=df.to_markdown(index=False),
            format="csv",
        )

    @staticmethod
    def chunk_text(text: str, size: int = CHUNK_SIZE) -> list[str]:
        """按 size 字符切片。"""
        if not text:
            return []
        return [text[i : i + size] for i in range(0, len(text), size)]

    @staticmethod
    def format_for_prompt(doc: ExtractedDocument, filename: str) -> str:
        """构造注入 LLM 上下文的附件文本。超长只注入前 2 片并提示 office_read。"""
        if len(doc.text) <= OfficeExtractor.INLINE_LIMIT:
            return (
                f"【用户上传附件：{filename}】\n"
                f"{doc.text or doc.markdown or '（未从文件中提取到文本内容，可能为纯图片扫描件）'}"
            )
        chunks = OfficeExtractor.chunk_text(doc.text)
        preview = "\n...\n".join(chunks[:2])
        return (
            f"【用户上传附件：{filename}】\n{preview}\n"
            f"\n（附件较长，以上仅展示前 {len(chunks[:2]) * OfficeExtractor.CHUNK_SIZE} 字符。"
            f"如需读取后续内容，请调用 office_read 工具按需读取。）"
        )
