import base64
import io

import pytest
from django.core.files.uploadedfile import SimpleUploadedFile

from smart_assistant.extractors.office_extractor import (
    OfficeExtractor,
    OfficeExtractError,
    ExtractedDocument,
)

# 合法最小 PDF（含文本 "Test PDF content"）。
# 注意：原 brief 版本 stream 字典缺 /Filter /FlateDecode（且 /Length 错），
# 导致 pdfplumber/pypdf 均解析失败；此版本补齐 Filter 并修正 Length 与 xref 偏移。
MIN_PDF_B64 = (
    "JVBERi0xLjQKMSAwIG9iago8PCAvVHlwZSAvQ2F0YWxvZyAvUGFnZXMgMiAwIFIgPj4KZW5kb2Jq"
    "CjIgMCBvYmoKPDwgL1R5cGUgL1BhZ2VzIC9LaWRzIFszIDAgUl0gL0NvdW50IDEgPj4KZW5kb2Jq"
    "CjMgMCBvYmoKPDwgL1R5cGUgL1BhZ2UgL1BhcmVudCAyIDAgUiAvTWVkaWFCb3ggWzAgMCA2MTIg"
    "NzkyXSAvQ29udGVudHMgNCAwIFIgL1Jlc291cmNlcyA8PCAvRm9udCA8PCAvRjEgNSAwIFIgPj4g"
    "Pj4gPj4KZW5kb2JqCjQgMCBvYmoKPDwgL0xlbmd0aCA1NSAvRmlsdGVyIC9GbGF0ZURlY29kZSA+"
    "PgpzdHJlYW0KeJxzClHQdzNUMDJRCElTMDdSMDcwUAhJUdAISS0uUQhwcVNIzs8rSc0r0VQIyVJw"
    "DQEAEtAMkgplbmRzdHJlYW0KZW5kb2JqCjUgMCBvYmoKPDwgL1R5cGUgL0ZvbnQgL1N1YnR5cGUg"
    "L1R5cGUxIC9CYXNlRm9udCAvSGVsdmV0aWNhID4+CmVuZG9iagp4cmVmCjAgNgowMDAwMDAwMDAw"
    "IDY1NTM1IGYgCjAwMDAwMDAwMDkgMDAwMDAgbiAKMDAwMDAwMDA1OCAwMDAwMCBuIAowMDAwMDAw"
    "MTE1IDAwMDAwIG4gCjAwMDAwMDAyNDEgMDAwMDAgbiAKMDAwMDAwMDM2NyAwMDAwMCBuIAp0cmFp"
    "bGVyCjw8IC9TaXplIDYgL1Jvb3QgMSAwIFIgPj4Kc3RhcnR4cmVmCjQzNwolJUVPRg=="
)


def _make_docx_bytes() -> bytes:
    from docx import Document

    buf = io.BytesIO()
    doc = Document()
    doc.add_paragraph("测试段落一")
    doc.add_paragraph("测试段落二")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "姓名"
    table.cell(0, 1).text = "部门"
    table.cell(1, 0).text = "张三"
    table.cell(1, 1).text = "技术部"
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    buf = io.BytesIO()
    wb = Workbook()
    ws = wb.active
    ws.title = "人员表"
    ws.append(["姓名", "部门", "人数"])
    ws.append(["张三", "技术部", 3])
    ws.append(["李四", "市场部", 5])
    wb.save(buf)
    return buf.getvalue()


def _make_pptx_bytes() -> bytes:
    from pptx import Presentation

    buf = io.BytesIO()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "项目汇报"
    body = slide.placeholders[1]
    body.text = "阶段一完成"
    prs.save(buf)
    return buf.getvalue()


def _upload(name: str, content: bytes) -> SimpleUploadedFile:
    return SimpleUploadedFile(name, content, content_type="application/octet-stream")


class TestOfficeExtractor:
    def test_docx_extracts_text_and_table(self):
        doc = OfficeExtractor.extract(_upload("合同.docx", _make_docx_bytes()))
        assert "测试段落一" in doc.text
        assert doc.tables and doc.tables[0]["headers"] == ["姓名", "部门"]
        assert "张三" in doc.markdown

    def test_pdf_extracts_text(self):
        doc = OfficeExtractor.extract(_upload("报告.pdf", base64.b64decode(MIN_PDF_B64)))
        assert "Test PDF content" in doc.text
        assert doc.format == "pdf"

    def test_pdf_extract_text_called_once_per_page(self):
        from unittest import mock

        page = mock.MagicMock()
        page.extract_text.return_value = "页面文字"
        page.extract_tables.return_value = []
        pdf = mock.MagicMock()
        pdf.__enter__.return_value = pdf
        pdf.pages = [page, page]
        with mock.patch("pdfplumber.open", return_value=pdf):
            doc = OfficeExtractor.extract(_upload("多页.pdf", base64.b64decode(MIN_PDF_B64)))
        assert doc.text == "页面文字\n页面文字"
        # 每页只调一次 extract_text（昂贵操作），不重复调用
        assert page.extract_text.call_count == 2

    def test_xlsx_extracts_sheets(self):
        doc = OfficeExtractor.extract(_upload("名单.xlsx", _make_xlsx_bytes()))
        assert len(doc.sheets) == 1
        assert doc.sheets[0]["name"] == "人员表"
        assert "市场部" in doc.text

    def test_pptx_extracts_text(self):
        doc = OfficeExtractor.extract(_upload("汇报.pptx", _make_pptx_bytes()))
        assert "项目汇报" in doc.text

    def test_txt_extracts_text(self):
        doc = OfficeExtractor.extract(_upload("笔记.txt", "纯文本内容".encode()))
        assert doc.text == "纯文本内容"

    def test_csv_gbk_encoding_fallback(self):
        content = "姓名,部门\n张三,技术部\n".encode("gbk")
        doc = OfficeExtractor.extract(_upload("名单.csv", content))
        assert doc.format == "csv"
        assert "张三" in doc.text
        assert "技术部" in doc.text

    def test_unsupported_extension_raises(self):
        with pytest.raises(OfficeExtractError):
            OfficeExtractor.extract(_upload("旧版.doc", b"\xd0\xcf\x11\xe0"))

    def test_corrupt_file_raises(self):
        with pytest.raises(OfficeExtractError):
            OfficeExtractor.extract(_upload("坏.docx", b"not a docx at all"))

    def test_oversized_file_rejected_before_read(self):
        class FakeOversizeFile:
            name = "超大.xlsx"
            size = OfficeExtractor.MAX_UPLOAD_SIZE + 1

            def __init__(self):
                self.read_called = False

            def read(self, *args, **kwargs):
                self.read_called = True
                return b"x" * 100

        f = FakeOversizeFile()
        with pytest.raises(OfficeExtractError):
            OfficeExtractor.extract(f)
        # file.size 预检应在 read() 之前拒绝，避免超大文件读入内存
        assert f.read_called is False

    def test_chunk_text_splits_by_size(self):
        chunks = OfficeExtractor.chunk_text("a" * 20_000, size=8_000)
        assert len(chunks) == 3
        assert all(len(c) <= 8_000 for c in chunks)

    def test_format_for_prompt_truncates_long_docs(self):
        # 注意：文件名用 .md 而非 .txt —— ".txt" 本身含一个 "x"，
        # 会使 prompt.count("x") 恒多 1（16001），导致断言 <= 16000 失败。
        long_doc = ExtractedDocument(text="x" * 120_000, format="txt")
        prompt = OfficeExtractor.format_for_prompt(long_doc, "长文.md")
        assert "长文.md" in prompt
        assert "office_read" in prompt
        assert prompt.count("x") <= 16_000
